import io
import uuid

import sentry_sdk
from fastapi import APIRouter, Depends, HTTPException
from pydantic import BaseModel

import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd
from langchain_text_splitters import RecursiveCharacterTextSplitter
from qdrant_client import models

from clients import supabase, qdrant, embeddings
from config import QDRANT_COLLECTION, MAX_CHUNKS_PER_INGEST
from auth import verify_token, require_project_access
from ratelimit import is_rate_limited
from usage import get_plan_limits

router = APIRouter()


def _purge_file_points(file_id: str):
    """Delete every Qdrant point belonging to a file."""
    qdrant.delete(
        collection_name=QDRANT_COLLECTION,
        points_selector=models.Filter(
            must=[models.FieldCondition(
                key="file_id",
                match=models.MatchValue(value=file_id)
            )]
        )
    )


# -------------------------------------------------
# MODELS
# -------------------------------------------------
class IngestRequest(BaseModel):
    projectId: str
    filename: str
    filePath: str


# -------------------------------------------------
# TEXT EXTRACTORS
# -------------------------------------------------
def extract_pdf(b):
    with pdfplumber.open(io.BytesIO(b)) as pdf:
        return [(i+1, p.extract_text() or "") for i, p in enumerate(pdf.pages) if p.extract_text()]

def extract_docx(b):
    d = Document(io.BytesIO(b))
    return [(1, "\n".join(p.text for p in d.paragraphs if p.text.strip()))]

def extract_pptx(b):
    prs = Presentation(io.BytesIO(b))
    out = []
    for i, s in enumerate(prs.slides):
        txt = "\n".join(sh.text for sh in s.shapes if hasattr(sh, "text"))
        if txt.strip():
            out.append((i+1, txt))
    return out

def extract_excel(b):
    xls = pd.ExcelFile(io.BytesIO(b))
    return [(n, xls.parse(n).astype(str).fillna("").to_csv(index=False)) for n in xls.sheet_names]

def extract_txt(b):
    return [(1, b.decode("utf-8", errors="ignore"))]


EXTRACTORS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "ppt": extract_pptx,
    "pptx": extract_pptx,
    "xls": extract_excel,
    "xlsx": extract_excel,
    "txt": extract_txt,
}


# -------------------------------------------------
# INGEST ENDPOINT
# -------------------------------------------------
@router.post("/ingest")
def ingest(req: IngestRequest, user=Depends(verify_token)):
    require_project_access(user.id, req.projectId, tab="documents")

    # Every ingest costs real OpenAI money, and nothing here was throttled
    # or capped before — a scripted loop could re-embed indefinitely.
    if is_rate_limited(f"ingest:{req.projectId}", limit=20, window_seconds=60):
        raise HTTPException(
            status_code=429,
            detail="Too many uploads at once. Please wait a minute and try again.",
        )

    # FIX: filePath is otherwise fully caller-controlled — without this
    # check, a user with a role on their OWN project could point filePath
    # at another project's storage object and have it indexed (crediting
    # req.projectId) as if it were their own document, exfiltrating another
    # tenant's file content into their own chatbot's knowledge base.
    if not req.filePath.startswith(f"{req.projectId}/"):
        raise HTTPException(status_code=403, detail="filePath does not belong to this project")

    row = supabase.table("files") \
        .select("id") \
        .eq("project_id", req.projectId) \
        .eq("filename", req.filename) \
        .execute()

    if not row.data:
        # Was returning this with HTTP 200, so the caller's .ok check passed
        # and a failed ingest looked identical to a successful one.
        raise HTTPException(status_code=404, detail="File not found")

    file_id = row.data[0]["id"]

    # Count OTHER documents in the project — this one's row already exists,
    # created by the upload route before it called us.
    limits = get_plan_limits(req.projectId)
    existing = supabase.table("files")         .select("id", count="exact")         .eq("project_id", req.projectId)         .neq("id", file_id)         .execute()
    if (existing.count or 0) >= limits["documents"]:
        # Clean up rather than leaving a failed row and a stored object the
        # user didn't get any value from.
        supabase.table("files").delete().eq("id", file_id).execute()
        try:
            supabase.storage.from_("documents").remove([req.filePath])
        except Exception:
            pass
        raise HTTPException(
            status_code=403,
            detail=f"You've reached your plan's limit of {limits['documents']} documents. Delete one, or upgrade your plan, to add more.",
        )

    supabase.table("files").update({"status": "processing"}).eq("id", file_id).execute()

    ext = req.filename.lower().split(".")[-1]
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        supabase.table("files").update({"status": "failed"}).eq("id", file_id).execute()
        raise HTTPException(status_code=400, detail="That file type isn't supported.")

    # Everything below can fail on someone else's infrastructure (Supabase
    # storage, OpenAI, Qdrant) or on a corrupt/password-protected file. Without
    # this, any of those left the row pinned at "processing" forever with no
    # reason recorded and an unhandled 500 to the caller.
    try:
        b = supabase.storage.from_("documents").download(req.filePath)

        splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
        chunks, metas = [], []

        for page, text in extractor(b):
            for c in splitter.split_text(text):
                chunks.append(c)
                metas.append({
                    "project_id": req.projectId,
                    "file_id": file_id,
                    "filename": req.filename,
                    "page_number": page,
                    "source_type": "document",
                    "text": c,
                })

        if not chunks:
            supabase.table("files").update({"status": "failed"}).eq("id", file_id).execute()
            raise HTTPException(
                status_code=400,
                detail="Couldn't read any text from that file. If it's a scanned PDF, it needs to contain selectable text.",
            )

        # A single enormous file would otherwise become one unbounded
        # embedding bill. Index the first N chunks and stop there.
        truncated = len(chunks) > MAX_CHUNKS_PER_INGEST
        if truncated:
            chunks = chunks[:MAX_CHUNKS_PER_INGEST]
            metas = metas[:MAX_CHUNKS_PER_INGEST]
            print(f"ingest truncated file {file_id} to {MAX_CHUNKS_PER_INGEST} chunks")

        vectors = embeddings.embed_documents(chunks)

        # Re-ingesting reuses the same file_id (the row is upserted), so
        # without this the previous version's chunks stayed in Qdrant
        # alongside the new ones and the bot kept answering from content
        # the user believed they had replaced.
        _purge_file_points(file_id)

        qdrant.upload_points(
            collection_name=QDRANT_COLLECTION,
            points=[
                models.PointStruct(
                    id=str(uuid.uuid4()),
                    vector=v,
                    payload=m
                ) for v, m in zip(vectors, metas)
            ]
        )
    except HTTPException:
        raise
    except Exception as e:
        sentry_sdk.capture_exception(e)
        print(f"ingest failed for file {file_id}: {e}")
        supabase.table("files").update({"status": "failed"}).eq("id", file_id).execute()
        raise HTTPException(
            status_code=502,
            detail="We couldn't process that file. Please try uploading it again.",
        )

    supabase.table("files").update({"status": "indexed"}).eq("id", file_id).execute()
    return {"status": "indexed", "chunks_indexed": len(chunks), "truncated": truncated}


@router.delete("/document/{file_id}")
def delete_document(file_id: str, user=Depends(verify_token)):
    row = supabase.table("files").select("project_id").eq("id", file_id).maybe_single().execute()
    file_row = row.data if row else None
    if not file_row:
        raise HTTPException(status_code=404, detail="Not found")
    require_project_access(user.id, file_row["project_id"], tab="documents")

    _purge_file_points(file_id)
    supabase.table("files").delete().eq("id", file_id).execute()
    return {"status": "deleted"}
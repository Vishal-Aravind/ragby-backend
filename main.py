import os
import io
import uuid
from datetime import datetime
from fastapi import Query, Depends, Security
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from typing import List, Optional

from fastapi import FastAPI, Request, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi import UploadFile, File, Form
from pydantic import BaseModel
from dotenv import load_dotenv

import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd

from supabase import create_client
from openai import OpenAI

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from fastapi.staticfiles import StaticFiles

import requests
from qdrant_client import QdrantClient, models

from sources.gsheets import sync_sheet
from sources.postgres import run_text_to_sql, introspect_schema, validate_url
from sources.excel import sync_excel_url, sync_excel_bytes
from sources.website import sync_website

import sqlalchemy
import hmac
import hashlib
import time

from starlette.responses import PlainTextResponse

import stripe

# -------------------------------------------------
# ENV & CLIENTS
# -------------------------------------------------
load_dotenv()

SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
QDRANT_URL = os.getenv("QDRANT_URL")
QDRANT_API_KEY = os.getenv("QDRANT_API_KEY")
QDRANT_COLLECTION = os.getenv("QDRANT_COLLECTION")
FRONTEND_URL = os.getenv("FRONTEND_URL", "http://localhost:3000")
VERIFY_TOKEN = os.getenv("VERIFY_TOKEN")

WHATSAPP_TOKEN = os.getenv("WHATSAPP_TOKEN")
WHATSAPP_PHONE_NUMBER_ID = os.getenv("WHATSAPP_PHONE_NUMBER_ID")
WHATSAPP_VERIFY_TOKEN = os.getenv("WHATSAPP_VERIFY_TOKEN")

SLACK_CLIENT_ID = os.getenv("SLACK_CLIENT_ID")
SLACK_CLIENT_SECRET = os.getenv("SLACK_CLIENT_SECRET")
SLACK_SIGNING_SECRET = os.getenv("SLACK_SIGNING_SECRET")

META_APP_ID = os.getenv("META_APP_ID")
META_APP_SECRET = os.getenv("META_APP_SECRET")

STRIPE_SECRET_KEY = os.getenv("STRIPE_SECRET_KEY")
STRIPE_WEBHOOK_SECRET = os.getenv("STRIPE_WEBHOOK_SECRET")
STRIPE_PRO_MONTHLY = os.getenv("STRIPE_PRO_MONTHLY")
STRIPE_PRO_YEARLY = os.getenv("STRIPE_PRO_YEARLY")
STRIPE_BUSINESS_MONTHLY = os.getenv("STRIPE_BUSINESS_MONTHLY")
STRIPE_BUSINESS_YEARLY = os.getenv("STRIPE_BUSINESS_YEARLY")
stripe.api_key = STRIPE_SECRET_KEY

PRICE_TO_PLAN = {
    STRIPE_PRO_MONTHLY: "pro",
    STRIPE_PRO_YEARLY: "pro",
    STRIPE_BUSINESS_MONTHLY: "business",
    STRIPE_BUSINESS_YEARLY: "business",
}

assert QDRANT_URL and QDRANT_API_KEY and QDRANT_COLLECTION

supabase = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)
qdrant = QdrantClient(url=QDRANT_URL, api_key=QDRANT_API_KEY)

embeddings = OpenAIEmbeddings(
    model="text-embedding-3-small",
    openai_api_key=OPENAI_API_KEY
)
openai_client = OpenAI(api_key=OPENAI_API_KEY)


# -------------------------------------------------
# APP
# -------------------------------------------------
app = FastAPI()
app.mount("/static", StaticFiles(directory="static"), name="static")
app.add_middleware(
    CORSMiddleware,
    allow_origins=[FRONTEND_URL],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# -------------------------------------------------
# AUTH
# -------------------------------------------------
bearer_scheme = HTTPBearer()

def verify_token(credentials: HTTPAuthorizationCredentials = Security(bearer_scheme)):
    token = credentials.credentials
    try:
        user_response = supabase.auth.get_user(token)
        if not user_response or not user_response.user:
            raise HTTPException(status_code=401, detail="Invalid or expired token")
        return user_response.user
    except Exception:
        raise HTTPException(status_code=401, detail="Invalid or expired token")


# -------------------------------------------------
# MODELS
# -------------------------------------------------
class IngestRequest(BaseModel):
    projectId: str
    filename: str
    filePath: str

class ChatMessage(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    projectId: str
    chatId: str
    message: str

class PublicChatRequest(BaseModel):
    projectId: str
    message: str
    sessionId: Optional[str] = None


# -------------------------------------------------
# PLAN LIMITS & RATE LIMITING
# -------------------------------------------------
PLAN_LIMITS = {
    "free":     {"conversations": 100,   "projects": 1},
    "pro":      {"conversations": 1000,  "projects": 5},
    "business": {"conversations": 7500,  "projects": None},
}

def get_current_month() -> str:
    return datetime.utcnow().strftime("%Y-%m")

def check_rate_limit(project_id: str) -> dict:
    proj = supabase.table("projects") \
        .select("user_id") \
        .eq("id", project_id) \
        .single() \
        .execute()

    if not proj.data:
        return {"allowed": False, "reason": "Project not found"}

    user_id = proj.data["user_id"]

    profile = supabase.table("profiles") \
        .select("plan") \
        .eq("id", user_id) \
        .single() \
        .execute()

    plan = "free"
    if profile.data and profile.data.get("plan"):
        plan = profile.data["plan"]

    limit = PLAN_LIMITS.get(plan, PLAN_LIMITS["free"])
    month = get_current_month()

    try:
        usage = supabase.table("usage") \
            .select("count") \
            .eq("user_id", user_id) \
            .eq("month", month) \
            .single() \
            .execute()
        current_count = usage.data.get("count", 0) if usage.data else 0
    except Exception:
        current_count = 0

    if current_count >= limit["conversations"]:
        return {
            "allowed": False,
            "reason": f"Monthly limit of {limit['conversations']} conversations reached. Please upgrade your plan.",
            "plan": plan,
            "usage": current_count,
            "limit": limit["conversations"],
        }

    return {
        "allowed": True,
        "plan": plan,
        "usage": current_count,
        "limit": limit["conversations"],
    }

def increment_usage(project_id: str):
    try:
        proj = supabase.table("projects") \
            .select("user_id") \
            .eq("id", project_id) \
            .single() \
            .execute()

        if not proj.data:
            return

        user_id = proj.data["user_id"]
        month = get_current_month()

        existing = supabase.table("usage") \
            .select("id, count") \
            .eq("user_id", user_id) \
            .eq("month", month) \
            .single() \
            .execute()

        if existing.data:
            supabase.table("usage") \
                .update({"count": existing.data["count"] + 1}) \
                .eq("id", existing.data["id"]) \
                .execute()
        else:
            supabase.table("usage").insert({
                "user_id": user_id,
                "month": month,
                "count": 1,
            }).execute()
    except Exception as e:
        print(f"increment_usage error: {e}")


# -------------------------------------------------
# USAGE STATUS ENDPOINT
# -------------------------------------------------
@app.get("/usage/status")
def usage_status(user=Depends(verify_token)):
    user_id = user.id
    month = get_current_month()

    profile = supabase.table("profiles") \
        .select("plan") \
        .eq("id", user_id) \
        .single() \
        .execute()

    plan = "free"
    if profile.data and profile.data.get("plan"):
        plan = profile.data["plan"]

    limit = PLAN_LIMITS.get(plan, PLAN_LIMITS["free"])

    try:
        usage = supabase.table("usage") \
            .select("count") \
            .eq("user_id", user_id) \
            .eq("month", month) \
            .single() \
            .execute()
        count = usage.data.get("count", 0) if usage.data else 0
    except Exception:
        count = 0

    return {
        "plan": plan,
        "usage": count,
        "limit": limit["conversations"],
        "remaining": max(0, limit["conversations"] - count),
        "month": month,
    }


# -------------------------------------------------
# TEXT EXTRACTORS
# -------------------------------------------------
def extract_pdf(b):
    with pdfplumber.open(io.BytesIO(b)) as pdf:
        return [(i+1, p.extract_text() or "") for i, p in enumerate(pdf.pages) if p.extract_text()]

def extract_docx(b):
    d = Document(io.BytesIO(b))
    return [(1, "\n".join(p.text for p in d.paragraphs if p.text.strip()))]

def extract_pptx(b):
    prs = Presentation(io.BytesIO(b))
    out = []
    for i, s in enumerate(prs.slides):
        txt = "\n".join(sh.text for sh in s.shapes if hasattr(sh, "text"))
        if txt.strip():
            out.append((i+1, txt))
    return out

def extract_excel(b):
    xls = pd.ExcelFile(io.BytesIO(b))
    return [(n, xls.parse(n).astype(str).fillna("").to_csv(index=False)) for n in xls.sheet_names]

def extract_txt(b):
    return [(1, b.decode("utf-8", errors="ignore"))]


# -------------------------------------------------
# INGEST
# -------------------------------------------------
@app.post("/ingest")
def ingest(req: IngestRequest, user=Depends(verify_token)):
    row = supabase.table("files").select("id").eq("project_id", req.projectId).eq("filename", req.filename).execute()
    if not row.data:
        return {"error": "file not found"}

    file_id = row.data[0]["id"]
    supabase.table("files").update({"status": "processing"}).eq("id", file_id).execute()

    b = supabase.storage.from_("documents").download(req.filePath)
    ext = req.filename.lower().split(".")[-1]

    extractor = {
        "pdf": extract_pdf,
        "docx": extract_docx,
        "ppt": extract_pptx,
        "pptx": extract_pptx,
        "xls": extract_excel,
        "xlsx": extract_excel,
        "txt": extract_txt,
    }.get(ext)

    if not extractor:
        supabase.table("files").update({"status": "failed"}).eq("id", file_id).execute()
        return {"error": "unsupported file type"}

    splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
    chunks, metas = [], []

    for page, text in extractor(b):
        for c in splitter.split_text(text):
            chunks.append(c)
            metas.append({
                "project_id": req.projectId,
                "file_id": file_id,
                "filename": req.filename,
                "page_number": page,
                "source_type": "document",
                "text": c,
            })

    vectors = embeddings.embed_documents(chunks)

    qdrant.upload_points(
        collection_name=QDRANT_COLLECTION,
        points=[
            models.PointStruct(
                id=str(uuid.uuid4()),
                vector=v,
                payload=m
            ) for v, m in zip(vectors, metas)
        ]
    )

    supabase.table("files").update({"status": "indexed"}).eq("id", file_id).execute()
    return {"status": "indexed", "chunks_indexed": len(chunks)}


# -------------------------------------------------
# CHAT HISTORY HELPERS
# -------------------------------------------------
def get_history(chat_id: str, limit: int = 7):
    res = supabase.table("chat_messages") \
        .select("role, content") \
        .eq("chat_id", chat_id) \
        .order("created_at", desc=True) \
        .limit(limit) \
        .execute()
    return list(reversed(res.data)) if res.data else []

def save_message(chat_id: str, role: str, content: str):
    supabase.table("chat_messages").insert({
        "chat_id": chat_id,
        "role": role,
        "content": content
    }).execute()


# -------------------------------------------------
# INTENT DETECTION
# -------------------------------------------------
def classify_intent(message: str) -> str:
    msg = message.lower().strip()
    words = msg.split()

    if len(words) <= 3 and msg in {"hi", "hello", "hey", "hi there", "hello there"}:
        return "greeting"

    if len(words) <= 4 and any(w in msg for w in ["thanks", "thank you", "thx"]):
        return "thanks"

    if any(k in msg for k in [
        "earlier", "previous", "you said", "we talked",
        "last message", "first question"
    ]):
        return "conversational"

    return "document_query"


def get_project_domain(project_id: str):
    res = supabase.table("projects") \
        .select("domain") \
        .eq("id", project_id) \
        .single() \
        .execute()
    return res.data.get("domain") if res.data else None


# -------------------------------------------------
# SOURCE INTENT CLASSIFIER
# -------------------------------------------------
def classify_source_intent(message: str) -> str:
    resp = openai_client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{
            "role": "user",
            "content": f"""Classify this question as either 'structured' or 'conceptual'.

'structured' = looking up a specific record, person, value, date, status, or list
Examples: "what are John's remarks", "show sales for March", "find order status for ID 123", "list all employees"

'conceptual' = asking about a process, policy, explanation, or general knowledge
Examples: "how does the refund process work", "what is the leave policy", "explain the onboarding steps"

Question: {message}
Reply with only one word: structured or conceptual"""
        }],
        temperature=0,
        max_tokens=10,
    )
    result = resp.choices[0].message.content.strip().lower()
    return result if result in ("structured", "conceptual") else "conceptual"


# -------------------------------------------------
# SYSTEM PROMPT
# -------------------------------------------------
SYSTEM_PROMPT = (
    "You are a helpful RAG AI assistant. Use ONLY the provided context and recent conversation if relevant. "
    "If the user asks about you, introduce yourself politely.\n\n"

    "Context source guidance:\n"
    "- Context chunks labeled [Source: gsheets] or [Source: database] contain structured data like records, names, values, dates.\n"
    "- Context chunks labeled [Source: document] contain policies, procedures, or explanatory content.\n"
    "- Always prefer the source that best matches the question type.\n\n"

    "Style:\n"
    "- Simple question → short answer\n"
    "- Complex question → structured answer, no heavy formatting, medium length\n"
    "- Be concise, do not over-explain unless necessary\n\n"
    "- Pricing/packages/plans questions → always list ALL options, never just one\n"

    "Logic:\n"
    "- One clear answer → answer directly\n"
    "- Multiple answers → ask a clarification question briefly\n"
    "- Partial info → answer + mention missing briefly\n"
    "- No answer → say you couldn't find the information and ask if they have more details\n"
    "- Match answer length to question complexity\n\n"

    "Formatting rules:\n"
    "- Use bullet points ONLY when needed\n\n"

    "Rules:\n"
    "- No hallucination\n"
    "- No external knowledge\n"
    "- Use previous conversation history for context when relevant"
)


# -------------------------------------------------
# CORE CHAT LOGIC
# -------------------------------------------------
def run_chat(project_id: str, chat_id: str, message: str, history: list):
    try:
        history = history or []
        domain = get_project_domain(project_id)

        system_prompt = SYSTEM_PROMPT
        if domain:
            system_prompt += f"\n\nDomain:\n- You are specialized in {domain}."

        intent = classify_intent(message)

        if intent == "greeting":
            return {"answer": "Hey! 👋 What can I help you with?", "sources": []}

        if intent == "thanks":
            return {"answer": "You're welcome! 😊", "sources": []}

        save_message(chat_id, "user", message)

        if intent == "conversational":
            messages = [{"role": "system", "content": system_prompt}]
            for h in history[-7:]:
                messages.append({"role": h["role"], "content": h["content"]})
            messages.append({"role": "user", "content": message})

            completion = openai_client.chat.completions.create(
                model="gpt-4o-mini",
                messages=messages,
                temperature=0.3,
                max_tokens=500,
            )
            answer = completion.choices[0].message.content.strip()
            save_message(chat_id, "assistant", answer)
            return {"answer": answer, "sources": []}

        source_intent = classify_source_intent(message)

        query_for_embedding = message
        if len(message.split()) <= 4 and history:
            last_user_msgs = [m for m in history if m["role"] == "user"]
            if last_user_msgs:
                query_for_embedding = last_user_msgs[-1]["content"] + " " + message

        q = embeddings.embed_query(query_for_embedding)
        context = None

        if source_intent == "structured":
            res = qdrant.query_points(
                collection_name=QDRANT_COLLECTION,
                query=q,
                limit=7,
                query_filter=models.Filter(
                    must=[
                        models.FieldCondition(key="project_id", match=models.MatchValue(value=project_id)),
                        models.FieldCondition(key="source_type", match=models.MatchAny(any=["gsheets", "excel"])),
                    ]
                )
            )
            hits = res.points

            if hits:
                context = "\n\n---\n\n".join(
                    f"[Source: gsheets]\n{h.payload.get('text', '')}" for h in hits
                )
            else:
                pg_source = supabase.table("data_sources") \
                    .select("config, allowed_schema") \
                    .eq("project_id", project_id) \
                    .eq("type", "postgres") \
                    .limit(1) \
                    .execute()

                if pg_source.data:
                    db_url = pg_source.data[0]["config"]["url"]
                    allowed_schema = pg_source.data[0].get("allowed_schema")
                    sql_result = run_text_to_sql(message, db_url, openai_client, allowed_schema)
                    context = f"[Source: database]\n{sql_result}"
                else:
                    source_intent = "conceptual"

        if source_intent == "conceptual":
            res = qdrant.query_points(
                collection_name=QDRANT_COLLECTION,
                query=q,
                limit=7,
                query_filter=models.Filter(
                    must=[
                        models.FieldCondition(key="project_id", match=models.MatchValue(value=project_id)),
                        models.FieldCondition(key="source_type", match=models.MatchAny(any=["document", "website"])),
                    ]
                )
            )
            hits = res.points

            if hits:
                context = "\n\n---\n\n".join(
                    f"[Source: document]\n{h.payload.get('text', '')}" for h in hits
                )

        if not context:
            answer = "I couldn't find that in your documents or data sources."
            save_message(chat_id, "assistant", answer)
            return {"answer": answer, "sources": []}

        messages = [{"role": "system", "content": system_prompt}]
        for h in history[-7:]:
            messages.append({"role": h["role"], "content": h["content"]})
        messages.append({
            "role": "user",
            "content": f"Context:\n{context}\n\nQuestion:\n{message}"
        })

        completion = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0.2,
            max_tokens=300,
        )

        answer = completion.choices[0].message.content.strip()
        save_message(chat_id, "assistant", answer)
        return {"answer": answer, "sources": []}

    except Exception as e:
        print(f"ERROR IN RUN_CHAT: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


# -------------------------------------------------
# CHAT ENDPOINTS
# -------------------------------------------------
@app.post("/chat")
def chat(req: ChatRequest, user=Depends(verify_token)):
    rate_check = check_rate_limit(req.projectId)
    if not rate_check["allowed"]:
        raise HTTPException(status_code=429, detail=rate_check["reason"])

    history = get_history(req.chatId, limit=7)
    result = run_chat(req.projectId, req.chatId, req.message, history)
    increment_usage(req.projectId)
    return result


@app.post("/public/chat")
def public_chat(req: PublicChatRequest):
    session_id = req.sessionId or str(uuid.uuid4())

    existing = supabase.table("chats").select("id").eq("id", session_id).execute()
    if not existing.data:
        supabase.table("chats").insert({
            "id": session_id,
            "project_id": req.projectId,
            "title": "Public Chat",
            "channel": "public",
        }).execute()

    rate_check = check_rate_limit(req.projectId)
    if not rate_check["allowed"]:
        return {
            "answer": "Sorry, this assistant has reached its monthly limit. Please try again next month.",
            "sessionId": session_id,
        }

    history = get_history(session_id, limit=7) if req.sessionId else []
    result = run_chat(req.projectId, session_id, req.message, history)
    result["sessionId"] = session_id
    increment_usage(req.projectId)
    return result


# -------------------------------------------------
# DELETE DOCUMENT
# -------------------------------------------------
@app.delete("/document/{file_id}")
def delete_document(file_id: str, user=Depends(verify_token)):
    qdrant.delete(
        collection_name=QDRANT_COLLECTION,
        points_selector=models.Filter(
            must=[models.FieldCondition(
                key="file_id",
                match=models.MatchValue(value=file_id)
            )]
        )
    )
    supabase.table("files").delete().eq("id", file_id).execute()
    return {"status": "deleted"}


@app.get("/health")
def health():
    return {"status": "ok"}


# -------------------------------------------------
# SOURCES
# -------------------------------------------------
@app.get("/sources")
def list_sources(project_id: str, user=Depends(verify_token)):
    res = supabase.table("data_sources") \
        .select("*") \
        .eq("project_id", project_id) \
        .execute()
    return res.data

@app.post("/sources/add")
def add_source(data: dict, user=Depends(verify_token)):
    res = supabase.table("data_sources").insert({
        "project_id": data["projectId"],
        "type": data["type"],
        "label": data.get("label") or data["type"],
        "config": data["config"],
        "allowed_schema": data.get("allowed_schema"),
    }).execute()
    source = res.data[0]

    skipped_tabs = []

    if data["type"] == "gsheets":
        cfg = data["config"]
        result = sync_sheet(
            cfg["sheet_id"], cfg["range"], data["projectId"],
            source["id"], qdrant, embeddings, QDRANT_COLLECTION
        )
        skipped_tabs = result.get("skipped_tabs", [])

    elif data["type"] == "excel_online":
        cfg = data["config"]
        sync_excel_url(
            cfg["url"], data["projectId"],
            source["id"], qdrant, embeddings, QDRANT_COLLECTION
        )

    elif data["type"] == "website":
        cfg = data["config"]
        result = sync_website(
            url=cfg["url"],
            project_id=data["projectId"],
            source_id=source["id"],
            qdrant=qdrant,
            embeddings=embeddings,
            collection=QDRANT_COLLECTION,
            full_site=cfg.get("full_site", True),
            max_pages=cfg.get("max_pages", 50),
        )
        if result["pages_indexed"] == 0:
            supabase.table("data_sources").delete().eq("id", source["id"]).execute()
            raise HTTPException(
                status_code=400,
                detail="Could not crawl this website. It may be blocking automated access."
            )

    return {"id": source["id"], "skipped_tabs": skipped_tabs}

@app.delete("/sources/{source_id}")
def delete_source(source_id: str, user=Depends(verify_token)):
    qdrant.delete(collection_name=QDRANT_COLLECTION, points_selector=models.Filter(
        must=[models.FieldCondition(key="source_id", match=models.MatchValue(value=source_id))]
    ))
    supabase.table("data_sources").delete().eq("id", source_id).execute()
    return {"status": "deleted"}

@app.post("/sources/sync/{source_id}")
def resync_source(source_id: str, user=Depends(verify_token)):
    res = supabase.table("data_sources").select("*").eq("id", source_id).single().execute()
    s = res.data

    qdrant.delete(
        collection_name=QDRANT_COLLECTION,
        points_selector=models.Filter(
            must=[models.FieldCondition(
                key="source_id",
                match=models.MatchValue(value=source_id)
            )]
        )
    )

    if s["type"] == "gsheets":
        cfg = s["config"]
        sync_sheet(
            cfg["sheet_id"], cfg["range"],
            s["project_id"], source_id,
            qdrant, embeddings, QDRANT_COLLECTION
        )
    elif s["type"] == "excel_online":
        cfg = s["config"]
        sync_excel_url(
            cfg["url"], s["project_id"],
            source_id, qdrant, embeddings, QDRANT_COLLECTION
        )
    elif s["type"] == "website":
        cfg = s["config"]
        sync_website(
            url=cfg["url"],
            project_id=s["project_id"],
            source_id=source_id,
            qdrant=qdrant,
            embeddings=embeddings,
            collection=QDRANT_COLLECTION,
            full_site=cfg.get("full_site", True),
            max_pages=cfg.get("max_pages", 50),
        )

    return {"status": "synced"}

@app.post("/sources/introspect")
def introspect(data: dict, user=Depends(verify_token)):
    db_url = data.get("db_url", "")
    try:
        validate_url(db_url)
        schema = introspect_schema(db_url)
        return {"schema": schema}
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not connect: {str(e)}")

@app.post("/sources/upload-excel")
async def upload_excel(
    file: UploadFile = File(...),
    projectId: str = Form(...),
    label: str = Form(""),
    source_id: str = Form(""),
    user=Depends(verify_token)
):
    file_bytes = await file.read()

    if source_id:
        sync_excel_bytes(file_bytes, projectId, source_id, qdrant, embeddings, QDRANT_COLLECTION)
        supabase.table("data_sources").update({
            "config": {"filename": file.filename},
            "label": label or file.filename,
        }).eq("id", source_id).execute()
        return {"id": source_id, "filename": file.filename}

    res = supabase.table("data_sources").insert({
        "project_id": projectId,
        "type": "excel_local",
        "label": label or file.filename,
        "config": {"filename": file.filename},
        "allowed_schema": None,
    }).execute()
    source = res.data[0]

    sync_excel_bytes(file_bytes, projectId, source["id"], qdrant, embeddings, QDRANT_COLLECTION)
    return {"id": source["id"], "filename": file.filename}


# -------------------------------------------------
# LEAD CAPTURE
# -------------------------------------------------
@app.get("/public/lead-config/{project_id}")
def get_lead_config(project_id: str):
    res = supabase.table("lead_capture_config") \
        .select("enabled, trigger_after_messages, form_title, form_subtitle") \
        .eq("project_id", project_id) \
        .execute()

    if not res.data or not res.data[0]["enabled"]:
        return {"enabled": False}

    return res.data[0]


class LeadSubmitRequest(BaseModel):
    project_id: str
    session_id: str
    name: str
    email: str
    phone: str

@app.post("/public/leads")
def submit_lead(req: LeadSubmitRequest):
    existing = supabase.table("leads") \
        .select("id") \
        .eq("session_id", req.session_id) \
        .eq("project_id", req.project_id) \
        .execute()

    if existing.data:
        return {"status": "already_captured"}

    if "@" not in req.email:
        raise HTTPException(status_code=400, detail="Invalid email")

    if len(req.phone.strip()) < 7:
        raise HTTPException(status_code=400, detail="Invalid phone")

    supabase.table("leads").insert({
        "project_id": req.project_id,
        "session_id": req.session_id,
        "name": req.name.strip(),
        "email": req.email.strip(),
        "phone": req.phone.strip(),
        "source": "widget",
    }).execute()

    return {"status": "captured"}


class LeadConfigRequest(BaseModel):
    projectId: str
    enabled: bool
    triggerAfterMessages: Optional[int] = 2
    formTitle: Optional[str] = "Before we continue..."
    formSubtitle: Optional[str] = "Please share your details to keep chatting."

@app.put("/lead-config")
def save_lead_config(req: LeadConfigRequest, user=Depends(verify_token)):
    supabase.table("lead_capture_config").upsert({
        "project_id": req.projectId,
        "enabled": req.enabled,
        "trigger_after_messages": req.triggerAfterMessages,
        "form_title": req.formTitle,
        "form_subtitle": req.formSubtitle,
    }, on_conflict="project_id").execute()
    return {"status": "saved"}


@app.get("/leads")
def get_leads(project_id: str, user=Depends(verify_token)):
    res = supabase.table("leads") \
        .select("*") \
        .eq("project_id", project_id) \
        .order("created_at", desc=True) \
        .execute()
    return res.data


# -------------------------------------------------
# TELEGRAM
# -------------------------------------------------
def send_telegram_message(bot_token: str, chat_id: int, text: str):
    url = f"https://api.telegram.org/bot{bot_token}/sendMessage"
    payload = {"chat_id": chat_id, "text": text, "parse_mode": "Markdown"}
    requests.post(url, json=payload)

def set_telegram_webhook(bot_token: str, webhook_url: str):
    url = f"https://api.telegram.org/bot{bot_token}/setWebhook"
    res = requests.post(url, json={"url": webhook_url})
    return res.json()

def get_bot_info(bot_token: str):
    url = f"https://api.telegram.org/bot{bot_token}/getMe"
    res = requests.get(url)
    return res.json()


@app.post("/telegram/connect")
def telegram_connect(data: dict, user=Depends(verify_token)):
    bot_token = data["bot_token"]
    project_id = data["projectId"]

    bot_info = get_bot_info(bot_token)
    if not bot_info.get("ok"):
        raise HTTPException(status_code=400, detail="Invalid bot token.")

    bot_username = bot_info["result"]["username"]

    supabase.table("telegram_integrations").upsert({
        "project_id": project_id,
        "bot_token": bot_token,
        "bot_username": bot_username,
    }, on_conflict="project_id").execute()

    webhook_url = f"{os.getenv('BACKEND_PUBLIC_URL')}/webhook/telegram/{project_id}"
    result = set_telegram_webhook(bot_token, webhook_url)

    if not result.get("ok"):
        raise HTTPException(status_code=400, detail=f"Could not set webhook: {result}")

    return {"success": True, "bot_username": bot_username}


@app.delete("/telegram/disconnect/{project_id}")
def telegram_disconnect(project_id: str, user=Depends(verify_token)):
    res = supabase.table("telegram_integrations") \
        .select("bot_token") \
        .eq("project_id", project_id) \
        .single() \
        .execute()

    if res.data:
        requests.post(f"https://api.telegram.org/bot{res.data['bot_token']}/deleteWebhook")

    supabase.table("telegram_integrations").delete().eq("project_id", project_id).execute()
    return {"success": True}


@app.get("/telegram/status/{project_id}")
def telegram_status(project_id: str, user=Depends(verify_token)):
    res = supabase.table("telegram_integrations") \
        .select("bot_username, created_at") \
        .eq("project_id", project_id) \
        .execute()

    if res.data:
        return {"connected": True, "bot_username": res.data[0]["bot_username"]}
    return {"connected": False}


@app.post("/webhook/telegram/{project_id}")
async def telegram_webhook(project_id: str, req: Request):
    body = await req.json()
    print(f"TELEGRAM BODY: {body}")

    try:
        message = body.get("message") or body.get("edited_message")
        if not message or "text" not in message:
            return {"status": "ignored"}

        text = message["text"]
        chat_id = message["chat"]["id"]
        chat_type = message["chat"]["type"]
        telegram_user_id = str(message["from"]["id"])
        username = message["from"].get("username") or message["from"].get("first_name", "User")

        res = supabase.table("telegram_integrations") \
            .select("bot_token, bot_username") \
            .eq("project_id", project_id) \
            .single() \
            .execute()

        if not res.data:
            return {"error": "integration not found"}

        bot_token = res.data["bot_token"]
        bot_username = res.data["bot_username"]

        if chat_type in ("group", "supergroup"):
            mention = f"@{bot_username}"
            if mention.lower() not in text.lower():
                return {"status": "ignored"}
            text = text.replace(mention, "").replace(mention.lower(), "").strip()
            if not text:
                send_telegram_message(bot_token, chat_id, "👋 Yes? Ask me anything!")
                return {"status": "ok"}

        if text.startswith("/start"):
            send_telegram_message(bot_token, chat_id, f"👋 Hi @{username}! I'm ready to help. Ask me anything!")
            return {"status": "ok"}

        if text.startswith("/"):
            return {"status": "ignored"}

        chat = supabase.table("chats") \
            .select("id") \
            .eq("project_id", project_id) \
            .eq("external_id", telegram_user_id) \
            .eq("channel", "telegram") \
            .limit(1) \
            .execute()

        if chat.data:
            chat_id_db = chat.data[0]["id"]
        else:
            new_chat = supabase.table("chats").insert({
                "project_id": project_id,
                "external_id": telegram_user_id,
                "channel": "telegram",
                "title": f"Telegram @{username}",
            }).execute()
            chat_id_db = new_chat.data[0]["id"]

        rate_check = check_rate_limit(project_id)
        if not rate_check["allowed"]:
            send_telegram_message(bot_token, chat_id, "⚠️ Monthly message limit reached. Please try again next month.")
            return {"status": "rate_limited"}

        history = get_history(chat_id_db, limit=5)
        result = run_chat(project_id, chat_id_db, text, history)
        send_telegram_message(bot_token, chat_id, result["answer"])
        increment_usage(project_id)
        return {"status": "ok"}

    except Exception as e:
        print(f"TELEGRAM WEBHOOK ERROR: {e}")
        return {"status": "error"}


# -------------------------------------------------
# SLACK
# -------------------------------------------------
def verify_slack_signature(body: bytes, timestamp: str, signature: str) -> bool:
    if abs(time.time() - int(timestamp)) > 300:
        return False
    sig_basestring = f"v0:{timestamp}:{body.decode('utf-8')}"
    my_sig = "v0=" + hmac.new(
        SLACK_SIGNING_SECRET.encode(),
        sig_basestring.encode(),
        hashlib.sha256
    ).hexdigest()
    return hmac.compare_digest(my_sig, signature)

def send_slack_message(access_token: str, channel: str, text: str):
    requests.post(
        "https://slack.com/api/chat.postMessage",
        headers={"Authorization": f"Bearer {access_token}"},
        json={"channel": channel, "text": text}
    )


@app.get("/slack/auth-url")
def slack_auth_url(project_id: str, user=Depends(verify_token)):
    redirect_uri = f"{FRONTEND_URL}/api/slack/callback"
    scopes = "app_mentions:read,chat:write,channels:history,im:history,im:write"
    url = (
        f"https://slack.com/oauth/v2/authorize"
        f"?client_id={SLACK_CLIENT_ID}"
        f"&scope={scopes}"
        f"&redirect_uri={redirect_uri}"
        f"&state={project_id}"
    )
    return {"url": url}


@app.post("/slack/callback")
def slack_callback(data: dict):
    code = data["code"]
    project_id = data["project_id"]
    redirect_uri = f"{FRONTEND_URL}/api/slack/callback"

    res = requests.post("https://slack.com/api/oauth.v2.access", data={
        "client_id": SLACK_CLIENT_ID,
        "client_secret": SLACK_CLIENT_SECRET,
        "code": code,
        "redirect_uri": redirect_uri,
    })
    token_data = res.json()

    if not token_data.get("ok"):
        raise HTTPException(status_code=400, detail=f"Slack OAuth failed: {token_data.get('error')}")

    supabase.table("slack_integrations").upsert({
        "project_id": project_id,
        "access_token": token_data["access_token"],
        "team_id": token_data["team"]["id"],
        "team_name": token_data["team"]["name"],
        "bot_user_id": token_data["bot_user_id"],
    }, on_conflict="project_id").execute()

    return {"success": True, "team_name": token_data["team"]["name"]}


@app.get("/slack/status/{project_id}")
def slack_status(project_id: str, user=Depends(verify_token)):
    res = supabase.table("slack_integrations") \
        .select("team_name, team_id") \
        .eq("project_id", project_id) \
        .execute()
    if res.data:
        return {"connected": True, "team_name": res.data[0]["team_name"]}
    return {"connected": False}


@app.delete("/slack/disconnect/{project_id}")
def slack_disconnect(project_id: str, user=Depends(verify_token)):
    supabase.table("slack_integrations").delete().eq("project_id", project_id).execute()
    return {"success": True}


@app.post("/webhook/slack")
async def slack_webhook(req: Request):
    body_bytes = await req.body()
    body = await req.json()

    if body.get("type") == "url_verification":
        return {"challenge": body["challenge"]}

    timestamp = req.headers.get("X-Slack-Request-Timestamp", "")
    signature = req.headers.get("X-Slack-Signature", "")
    if not verify_slack_signature(body_bytes, timestamp, signature):
        raise HTTPException(status_code=403, detail="Invalid signature")

    event = body.get("event", {})
    event_type = event.get("type")

    if event_type not in ("app_mention", "message"):
        return {"status": "ignored"}

    if event.get("bot_id") or event.get("subtype"):
        return {"status": "ignored"}

    text = event.get("text", "").strip()
    channel = event.get("channel")
    user_id = event.get("user")
    team_id = body.get("team_id")

    if not text or not channel or not user_id:
        return {"status": "ignored"}

    res = supabase.table("slack_integrations") \
        .select("project_id, access_token, bot_user_id") \
        .eq("team_id", team_id) \
        .single() \
        .execute()

    if not res.data:
        return {"error": "integration not found"}

    project_id = res.data["project_id"]
    access_token = res.data["access_token"]
    bot_user_id = res.data["bot_user_id"]

    text = text.replace(f"<@{bot_user_id}>", "").strip()
    if not text:
        send_slack_message(access_token, channel, "👋 Yes? Ask me anything!")
        return {"status": "ok"}

    chat = supabase.table("chats") \
        .select("id") \
        .eq("project_id", project_id) \
        .eq("external_id", user_id) \
        .eq("channel", "slack") \
        .limit(1) \
        .execute()

    if chat.data:
        chat_id = chat.data[0]["id"]
    else:
        new_chat = supabase.table("chats").insert({
            "project_id": project_id,
            "external_id": user_id,
            "channel": "slack",
            "title": f"Slack {user_id}",
        }).execute()
        chat_id = new_chat.data[0]["id"]

    rate_check = check_rate_limit(project_id)
    if not rate_check["allowed"]:
        send_slack_message(access_token, channel, "⚠️ Monthly message limit reached. Please try again next month.")
        return {"status": "rate_limited"}

    history = get_history(chat_id, limit=5)
    result = run_chat(project_id, chat_id, text, history)
    send_slack_message(access_token, channel, result["answer"])
    increment_usage(project_id)
    return {"status": "ok"}


# -------------------------------------------------
# WHATSAPP
# -------------------------------------------------
def send_whatsapp_message(to: str, text: str, phone_number_id: str = None, token: str = None):
    pid = phone_number_id or WHATSAPP_PHONE_NUMBER_ID
    tok = token or WHATSAPP_TOKEN
    url = f"https://graph.facebook.com/v19.0/{pid}/messages"
    headers = {
        "Authorization": f"Bearer {tok}",
        "Content-Type": "application/json",
    }
    payload = {
        "messaging_product": "whatsapp",
        "to": to,
        "type": "text",
        "text": {"body": text},
    }
    res = requests.post(url, headers=headers, json=payload)
    if not res.ok:
        print(f"WhatsApp send error: {res.text}")
    return res


@app.get("/webhook/whatsapp")
async def whatsapp_verify(request: Request):
    params = dict(request.query_params)
    mode = params.get("hub.mode")
    token = params.get("hub.verify_token")
    challenge = params.get("hub.challenge")

    if mode == "subscribe" and token == WHATSAPP_VERIFY_TOKEN:
        print("WhatsApp webhook verified")
        return PlainTextResponse(challenge)

    raise HTTPException(status_code=403, detail="Verification failed")


@app.post("/webhook/whatsapp")
async def whatsapp_webhook(request: Request):
    body = await request.json()

    try:
        entry = body.get("entry", [])[0]
        changes = entry.get("changes", [])[0]
        value = changes.get("value", {})

        if "statuses" in value:
            return {"status": "ignored"}

        messages = value.get("messages", [])
        if not messages:
            return {"status": "ignored"}

        message = messages[0]
        msg_type = message.get("type")

        if msg_type != "text":
            return {"status": "ignored"}

        from_number = message["from"]
        text = message["text"]["body"].strip()
        waba_id = value.get("metadata", {}).get("phone_number_id")

        res = supabase.table("whatsapp_integrations") \
            .select("project_id") \
            .eq("phone_number_id", waba_id or WHATSAPP_PHONE_NUMBER_ID) \
            .single() \
            .execute()

        if not res.data:
            print(f"No project found for phone_number_id: {waba_id}")
            return {"status": "ignored"}

        project_id = res.data["project_id"]

        chat = supabase.table("chats") \
            .select("id") \
            .eq("project_id", project_id) \
            .eq("external_id", from_number) \
            .eq("channel", "whatsapp") \
            .limit(1) \
            .execute()

        if chat.data:
            chat_id = chat.data[0]["id"]
        else:
            new_chat = supabase.table("chats").insert({
                "project_id": project_id,
                "external_id": from_number,
                "channel": "whatsapp",
                "title": f"WhatsApp {from_number}",
            }).execute()
            chat_id = new_chat.data[0]["id"]

        rate_check = check_rate_limit(project_id)
        if not rate_check["allowed"]:
            send_whatsapp_message(from_number, "⚠️ Monthly message limit reached. Please try again next month.")
            return {"status": "rate_limited"}

        history = get_history(chat_id, limit=5)
        result = run_chat(project_id, chat_id, text, history)
        send_whatsapp_message(from_number, result["answer"])
        increment_usage(project_id)
        return {"status": "ok"}

    except Exception as e:
        print(f"WHATSAPP WEBHOOK ERROR: {e}")
        return {"status": "error"}


@app.get("/whatsapp/status/{project_id}")
def whatsapp_status(project_id: str, user=Depends(verify_token)):
    res = supabase.table("whatsapp_integrations") \
        .select("phone_number_id, display_phone_number, waba_id") \
        .eq("project_id", project_id) \
        .execute()

    if res.data:
        return {"connected": True, **res.data[0]}
    return {"connected": False}


@app.post("/whatsapp/connect")
def whatsapp_connect(data: dict, user=Depends(verify_token)):
    project_id = data["projectId"]
    phone_number_id = data["phone_number_id"]
    waba_id = data.get("waba_id", "")
    display_phone_number = data.get("display_phone_number", "")

    supabase.table("whatsapp_integrations").upsert({
        "project_id": project_id,
        "phone_number_id": phone_number_id,
        "waba_id": waba_id,
        "display_phone_number": display_phone_number,
    }, on_conflict="project_id").execute()

    return {"success": True}


@app.delete("/whatsapp/disconnect/{project_id}")
def whatsapp_disconnect(project_id: str, user=Depends(verify_token)):
    supabase.table("whatsapp_integrations") \
        .delete() \
        .eq("project_id", project_id) \
        .execute()
    return {"success": True}


@app.post("/whatsapp/onboard")
def whatsapp_onboard(data: dict, user=Depends(verify_token)):
    code = data["code"]
    project_id = data["projectId"]

    token_res = requests.get(
        "https://graph.facebook.com/v19.0/oauth/access_token",
        params={
            "client_id": META_APP_ID,
            "client_secret": META_APP_SECRET,
            "code": code,
        }
    )
    token_data = token_res.json()
    print(f"Token exchange: {token_data}")

    if "access_token" not in token_data:
        raise HTTPException(status_code=400, detail=f"Token exchange failed: {token_data}")

    access_token = token_data["access_token"]

    waba_res = requests.get(
        "https://graph.facebook.com/v19.0/me/whatsapp_business_accounts",
        params={"access_token": access_token}
    )
    waba_data = waba_res.json()
    print(f"WABA data: {waba_data}")

    waba_id = waba_data.get("data", [{}])[0].get("id", "")

    phone_res = requests.get(
        f"https://graph.facebook.com/v19.0/{waba_id}/phone_numbers",
        params={"access_token": access_token}
    )
    phone_data = phone_res.json()
    print(f"Phone data: {phone_data}")

    phone_number_id = phone_data.get("data", [{}])[0].get("id", "")
    display_phone = phone_data.get("data", [{}])[0].get("display_phone_number", "")

    supabase.table("whatsapp_integrations").upsert({
        "project_id": project_id,
        "phone_number_id": phone_number_id,
        "waba_id": waba_id,
        "display_phone_number": display_phone,
    }, on_conflict="project_id").execute()

    return {
        "success": True,
        "phone_number_id": phone_number_id,
        "display_phone_number": display_phone,
        "waba_id": waba_id,
    }


# ─────────────────────────────────────────────────────────
# CREATE CHECKOUT SESSION
# ─────────────────────────────────────────────────────────
@app.post("/stripe/checkout")
def create_checkout(data: dict, user=Depends(verify_token)):
    price_id = data["priceId"]
    user_id = user.id
    user_email = user.email
 
    # Get or create Stripe customer
    profile = supabase.table("profiles") \
        .select("stripe_customer_id") \
        .eq("id", user_id) \
        .single() \
        .execute()
 
    stripe_customer_id = None
    if profile.data:
        stripe_customer_id = profile.data.get("stripe_customer_id")
 
    if not stripe_customer_id:
        customer = stripe.Customer.create(
            email=user_email,
            metadata={"supabase_user_id": user_id}
        )
        stripe_customer_id = customer.id
        supabase.table("profiles").upsert({
            "id": user_id,
            "stripe_customer_id": stripe_customer_id,
        }, on_conflict="id").execute()
 
    session = stripe.checkout.Session.create(
        customer=stripe_customer_id,
        payment_method_types=["card"],
        line_items=[{"price": price_id, "quantity": 1}],
        mode="subscription",
        success_url=f"{FRONTEND_URL}/dashboard?upgraded=true",
        cancel_url=f"{FRONTEND_URL}/pricing?cancelled=true",
        metadata={
            "supabase_user_id": user_id,
            "price_id": price_id,
        },
        subscription_data={
            "metadata": {
                "supabase_user_id": user_id,
                "price_id": price_id,
            }
        }
    )
 
    return {"url": session.url}
 
 
# ─────────────────────────────────────────────────────────
# CUSTOMER PORTAL — manage subscription
# ─────────────────────────────────────────────────────────
@app.post("/stripe/portal")
def customer_portal(user=Depends(verify_token)):
    user_id = user.id
 
    profile = supabase.table("profiles") \
        .select("stripe_customer_id") \
        .eq("id", user_id) \
        .single() \
        .execute()
 
    if not profile.data or not profile.data.get("stripe_customer_id"):
        raise HTTPException(status_code=400, detail="No Stripe customer found.")
 
    session = stripe.billing_portal.Session.create(
        customer=profile.data["stripe_customer_id"],
        return_url=f"{FRONTEND_URL}/dashboard",
    )
 
    return {"url": session.url}
 
 
# ─────────────────────────────────────────────────────────
# STRIPE WEBHOOK — handle subscription events
# ─────────────────────────────────────────────────────────
@app.post("/webhook/stripe")
async def stripe_webhook(request: Request):
    payload = await request.body()
    sig_header = request.headers.get("stripe-signature")
 
    try:
        event = stripe.Webhook.construct_event(
            payload, sig_header, STRIPE_WEBHOOK_SECRET
        )
    except stripe.error.SignatureVerificationError:
        raise HTTPException(status_code=400, detail="Invalid signature")
 
    event_type = event["type"]
    data = event["data"]["object"]
 
    # ── Checkout completed → activate plan ──────────────
    if event_type == "checkout.session.completed":
        user_id = data.get("metadata", {}).get("supabase_user_id")
        price_id = data.get("metadata", {}).get("price_id")
        subscription_id = data.get("subscription")
 
        if user_id and price_id:
            plan = PRICE_TO_PLAN.get(price_id, "free")
            supabase.table("profiles").upsert({
                "id": user_id,
                "plan": plan,
                "stripe_subscription_id": subscription_id,
            }, on_conflict="id").execute()
            print(f"Plan activated: {user_id} → {plan}")
 
    # ── Subscription updated → update plan ──────────────
    elif event_type == "customer.subscription.updated":
        subscription_id = data.get("id")
        price_id = data["items"]["data"][0]["price"]["id"]
        status = data.get("status")
 
        profile = supabase.table("profiles") \
            .select("id") \
            .eq("stripe_subscription_id", subscription_id) \
            .single() \
            .execute()
 
        if profile.data:
            user_id = profile.data["id"]
            if status in ("active", "trialing"):
                plan = PRICE_TO_PLAN.get(price_id, "free")
            else:
                plan = "free"
 
            supabase.table("profiles").update({
                "plan": plan,
            }).eq("id", user_id).execute()
            print(f"Plan updated: {user_id} → {plan} (status: {status})")
 
    # ── Subscription deleted → downgrade to free ────────
    elif event_type == "customer.subscription.deleted":
        subscription_id = data.get("id")
 
        profile = supabase.table("profiles") \
            .select("id") \
            .eq("stripe_subscription_id", subscription_id) \
            .single() \
            .execute()
 
        if profile.data:
            supabase.table("profiles").update({
                "plan": "free",
                "stripe_subscription_id": None,
            }).eq("id", profile.data["id"]).execute()
            print(f"Plan cancelled: {profile.data['id']} → free")
 
    return {"status": "ok"}
 
 
# ─────────────────────────────────────────────────────────
# GET CURRENT PLAN
# ─────────────────────────────────────────────────────────
@app.get("/stripe/plan")
def get_plan(user=Depends(verify_token)):
    user_id = user.id
 
    profile = supabase.table("profiles") \
        .select("plan, stripe_customer_id, stripe_subscription_id") \
        .eq("id", user_id) \
        .single() \
        .execute()
 
    if not profile.data:
        return {"plan": "free"}
 
    return {
        "plan": profile.data.get("plan", "free"),
        "has_subscription": bool(profile.data.get("stripe_subscription_id")),
    }
 